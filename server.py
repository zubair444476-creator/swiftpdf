import os
import io
import zipfile
import uuid
import base64

from flask import Flask, request, send_file, jsonify, render_template, send_from_directory
from werkzeug.utils import secure_filename

import fitz  # pymupdf
from PIL import Image

from docx import Document
from docx.shared import Pt

from pptx import Presentation
from pptx.util import Emu

from openpyxl import Workbook, load_workbook

from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.pdfgen import canvas as pdf_canvas

app = Flask(__name__)

MAX_CONTENT_LENGTH = 50 * 1024 * 1024  # 50 MB
app.config["MAX_CONTENT_LENGTH"] = MAX_CONTENT_LENGTH


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def get_uploaded_files(field_name="files"):
    files = request.files.getlist(field_name)
    if not files:
        single = request.files.get("file")
        if single:
            files = [single]
    return [f for f in files if f and f.filename]


def send_bytes(data: bytes, filename: str, mimetype: str):
    return send_file(
        io.BytesIO(data),
        mimetype=mimetype,
        as_attachment=True,
        download_name=filename,
    )


def parse_page_spec(spec: str, page_count: int):
    """Parse a 1-indexed page spec like '1,3,5-7' into a 0-indexed list,
    preserving the order given (so it can also be used to reorder pages)."""
    result = []
    spec = (spec or "").strip()
    if not spec:
        return list(range(page_count))

    for chunk in spec.split(","):
        chunk = chunk.strip()
        if not chunk:
            continue
        if "-" in chunk:
            start_s, end_s = chunk.split("-", 1)
            start, end = int(start_s), int(end_s)
            step = 1 if end >= start else -1
            for p in range(start, end + step, step):
                if 1 <= p <= page_count:
                    result.append(p - 1)
        else:
            p = int(chunk)
            if 1 <= p <= page_count:
                result.append(p - 1)
    return result


# ---------------------------------------------------------------------------
# Static frontend
# ---------------------------------------------------------------------------

@app.route("/")
def index():
    return send_from_directory(".", "index.html")


@app.route("/health")
def health():
    return jsonify({"status": "ok"})


# ---------------------------------------------------------------------------
# Page thumbnails (used by the visual page picker in the frontend)
# ---------------------------------------------------------------------------

@app.route("/api/page-thumbnails", methods=["POST"])
def page_thumbnails():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a PDF file."}), 400

    src = files[0]
    data = src.read()

    zoom = 0.4          # ~58 DPI — small enough to be fast, big enough to read
    mat = fitz.Matrix(zoom, zoom)
    thumbs = []

    with fitz.open(stream=data, filetype="pdf") as doc:
        page_count = len(doc)
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat)
            png_bytes = pix.tobytes("png")
            b64 = base64.b64encode(png_bytes).decode("ascii")
            thumbs.append({
                "page": i + 1,           # 1-indexed
                "width": pix.width,
                "height": pix.height,
                "dataUrl": f"data:image/png;base64,{b64}",
            })

    return jsonify({"pageCount": page_count, "thumbnails": thumbs})


# ---------------------------------------------------------------------------
# Merge PDFs
# ---------------------------------------------------------------------------

@app.route("/api/merge", methods=["POST"])
def merge_pdfs():
    files = get_uploaded_files()
    if len(files) < 2:
        return jsonify({"error": "Upload at least 2 PDF files to merge."}), 400

    merged = fitz.open()
    try:
        for f in files:
            data = f.read()
            with fitz.open(stream=data, filetype="pdf") as doc:
                merged.insert_pdf(doc)

        out_bytes = merged.tobytes()
    finally:
        merged.close()

    return send_bytes(out_bytes, "merged.pdf", "application/pdf")


# ---------------------------------------------------------------------------
# Split a PDF into individual pages (returned as a zip)
# ---------------------------------------------------------------------------

@app.route("/api/split", methods=["POST"])
def split_pdf():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a PDF file to split."}), 400

    src = files[0]
    data = src.read()
    base_name = os.path.splitext(secure_filename(src.filename))[0] or "page"

    zip_buffer = io.BytesIO()
    with fitz.open(stream=data, filetype="pdf") as doc:
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            for i in range(len(doc)):
                single = fitz.open()
                single.insert_pdf(doc, from_page=i, to_page=i)
                page_bytes = single.tobytes()
                single.close()
                zf.writestr(f"{base_name}_page_{i + 1}.pdf", page_bytes)

    zip_buffer.seek(0)
    return send_bytes(zip_buffer.getvalue(), f"{base_name}_split.zip", "application/zip")


# ---------------------------------------------------------------------------
# Rotate a PDF (all pages) by 90/180/270 degrees
# ---------------------------------------------------------------------------

@app.route("/api/rotate", methods=["POST"])
def rotate_pdf():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a PDF file to rotate."}), 400

    try:
        angle = int(request.form.get("angle", 90))
    except ValueError:
        angle = 90
    angle = angle % 360

    src = files[0]
    data = src.read()

    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            page.set_rotation((page.rotation + angle) % 360)
        out_bytes = doc.tobytes()

    return send_bytes(out_bytes, "rotated.pdf", "application/pdf")


# ---------------------------------------------------------------------------
# Compress a PDF (re-save with garbage collection + deflate, downsample images)
# ---------------------------------------------------------------------------

@app.route("/api/compress", methods=["POST"])
def compress_pdf():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a PDF file to compress."}), 400

    quality = request.form.get("quality", "recommended")
    jpeg_quality = {"low": 40, "recommended": 60, "high": 80}.get(quality, 60)
    max_dim = {"low": 800, "recommended": 1200, "high": 1600}.get(quality, 1200)

    src = files[0]
    data = src.read()

    doc = fitz.open(stream=data, filetype="pdf")
    try:
        for page in doc:
            for img in page.get_images(full=True):
                xref = img[0]
                try:
                    base = doc.extract_image(xref)
                    img_bytes = base["image"]
                    pil_img = Image.open(io.BytesIO(img_bytes)).convert("RGB")

                    w, h = pil_img.size
                    scale = min(1.0, max_dim / max(w, h))
                    if scale < 1.0:
                        pil_img = pil_img.resize(
                            (max(1, int(w * scale)), max(1, int(h * scale))),
                            Image.LANCZOS,
                        )

                    buf = io.BytesIO()
                    pil_img.save(buf, format="JPEG", quality=jpeg_quality, optimize=True)
                    doc.update_stream(xref, buf.getvalue())
                except Exception:
                    # Not all images can be re-encoded (masks, CMYK, etc.) - skip those.
                    continue

        out_bytes = doc.tobytes(deflate=True, garbage=4)
    finally:
        doc.close()

    return send_bytes(out_bytes, "compressed.pdf", "application/pdf")


# ---------------------------------------------------------------------------
# Images -> PDF
# ---------------------------------------------------------------------------

@app.route("/api/images-to-pdf", methods=["POST"])
def images_to_pdf():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload one or more images."}), 400

    pil_images = []
    for f in files:
        img = Image.open(f.stream).convert("RGB")
        pil_images.append(img)

    buf = io.BytesIO()
    first, rest = pil_images[0], pil_images[1:]
    first.save(buf, format="PDF", save_all=True, append_images=rest)
    buf.seek(0)

    return send_bytes(buf.getvalue(), "images.pdf", "application/pdf")


# ---------------------------------------------------------------------------
# PDF -> Images (returned as a zip of PNGs)
# ---------------------------------------------------------------------------

@app.route("/api/pdf-to-images", methods=["POST"])
def pdf_to_images():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a PDF file to convert."}), 400

    src = files[0]
    data = src.read()
    base_name = os.path.splitext(secure_filename(src.filename))[0] or "page"

    zoom = 2.0  # ~144 DPI
    mat = fitz.Matrix(zoom, zoom)

    zip_buffer = io.BytesIO()
    with fitz.open(stream=data, filetype="pdf") as doc:
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=mat)
                zf.writestr(f"{base_name}_page_{i + 1}.png", pix.tobytes("png"))

    zip_buffer.seek(0)
    return send_bytes(zip_buffer.getvalue(), f"{base_name}_images.zip", "application/zip")


# ---------------------------------------------------------------------------
# PDF -> plain text (simple extraction)
# ---------------------------------------------------------------------------

@app.route("/api/pdf-to-text", methods=["POST"])
def pdf_to_text():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a PDF file to convert."}), 400

    src = files[0]
    data = src.read()
    base_name = os.path.splitext(secure_filename(src.filename))[0] or "document"

    text_parts = []
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            text_parts.append(page.get_text())

    full_text = "\n\n".join(text_parts).encode("utf-8")
    return send_bytes(full_text, f"{base_name}.txt", "text/plain")


# ---------------------------------------------------------------------------
# PDF -> Word (.docx)  — rich conversion preserving headings, tables, formatting
# ---------------------------------------------------------------------------

# This is the replacement for the pdf_to_word route in server.py
# Replace everything from line 335 to line 654 with this content

def _pdf_rgb(color_int):
    """PyMuPDF integer color → docx RGBColor."""
    from docx.shared import RGBColor
    return RGBColor((color_int >> 16) & 0xFF, (color_int >> 8) & 0xFF, color_int & 0xFF)


def _set_cell_shading(cell, hex_fill):
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn("w:shd")):
        tcPr.remove(old)
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_fill)
    tcPr.append(shd)


def _set_table_borders(table, color="B0C4D8", sz="4"):
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    tbl = table._tbl
    tblPr = tbl.tblPr
    borders = OxmlElement("w:tblBorders")
    for side in ("top", "left", "bottom", "right", "insideH", "insideV"):
        el = OxmlElement(f"w:{side}")
        el.set(qn("w:val"), "single")
        el.set(qn("w:sz"), sz)
        el.set(qn("w:color"), color)
        borders.append(el)
    tblPr.append(borders)


def _set_cell_border(cell, color="B0C4D8", sz="4"):
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcBorders = OxmlElement("w:tcBorders")
    for side in ("top", "left", "bottom", "right"):
        el = OxmlElement(f"w:{side}")
        el.set(qn("w:val"), "single")
        el.set(qn("w:sz"), sz)
        el.set(qn("w:color"), color)
        tcBorders.append(el)
    tcPr.append(tcBorders)


def _add_image_para(document, img_bytes, width_cm):
    """Add an image as a paragraph with no spacing."""
    from docx.shared import Cm, Pt
    p = document.add_paragraph()
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after = Pt(0)
    run = p.add_run()
    run.add_picture(io.BytesIO(img_bytes), width=Cm(width_cm))
    return p


def _clip_page_region(page, rect, zoom=3.0):
    """Render a rectangular region of a PDF page as PNG bytes."""
    mat = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=mat, clip=rect)
    return pix.tobytes("png")


def _detect_table_regions(page):
    """
    Find table bounding boxes on the page using PyMuPDF's find_tables.
    Returns list of fitz.Rect.
    """
    table_rects = []
    try:
        found = page.find_tables()
        for ft in found.tables:
            raw = ft.extract()
            ncols = max(len(r) for r in raw) if raw else 0
            flat = [c for r in raw for c in r if c and str(c).strip()]
            w = ft.bbox[2] - ft.bbox[0]
            if len(raw) >= 2 and ncols >= 2 and len(flat) >= 2 and w < page.rect.width - 5:
                table_rects.append(fitz.Rect(ft.bbox))
    except Exception:
        pass
    return table_rects


def _extract_table_data(page, table_rect):
    """
    Extract table content from a page region.
    Returns list of rows, each row is a list of cell text strings.
    Also detects header row background color.
    """
    blocks = page.get_text("dict")["blocks"]
    y_min, y_max = table_rect.y0, table_rect.y1
    x_min, x_max = table_rect.x0, table_rect.x1

    text_blocks = [
        b for b in blocks
        if b.get("type") == 0
        and (y_min - 5) <= b["bbox"][1] <= (y_max + 5)
        and b["bbox"][0] >= (x_min - 8)
        and b["bbox"][2] <= (x_max + 8)
    ]
    if not text_blocks:
        return [], None, None

    text_blocks.sort(key=lambda b: (round(b["bbox"][1] / 4) * 4, b["bbox"][0]))

    # Detect column boundaries
    x_lefts = sorted(set(round(b["bbox"][0]) for b in text_blocks))
    GAP = 15
    clusters = [[x_lefts[0]]]
    for x in x_lefts[1:]:
        if x - clusters[-1][-1] > GAP:
            clusters.append([])
        clusters[-1].append(x)
    ncols = len(clusters)
    boundaries = [(clusters[i][-1] + clusters[i+1][0]) / 2.0 for i in range(len(clusters)-1)]

    def col_index(x0):
        for i, bnd in enumerate(boundaries):
            if x0 < bnd:
                return i
        return ncols - 1

    y_groups = {}
    for b in text_blocks:
        y_key = round(b["bbox"][1] / 4) * 4
        y_groups.setdefault(y_key, []).append(b)

    rows = []
    first_row_spans = None
    for yi, y_key in enumerate(sorted(y_groups.keys())):
        cols_text = [""] * ncols
        cols_spans = [[] for _ in range(ncols)]
        for b in y_groups[y_key]:
            ci = col_index(b["bbox"][0])
            for line in b["lines"]:
                spans = [s for s in line["spans"] if s["text"].strip()]
                if not spans:
                    continue
                txt = "".join(s["text"] for s in spans).strip()
                sep = " " if cols_text[ci] else ""
                cols_text[ci] = cols_text[ci] + sep + txt
                cols_spans[ci].extend(spans)
        if not any(c.strip() for c in cols_text):
            continue
        if yi == 0:
            first_row_spans = cols_spans
        rows.append(cols_text)

    # Detect header fill color
    hdr_fill = "2F5496"
    try:
        row_top = table_rect.y0
        candidates = []
        for d in page.get_drawings():
            r = d.get("rect")
            if r is None:
                continue
            f = d.get("fill")
            if not f or f in ((1.0, 1.0, 1.0), (0.0, 0.0, 0.0)):
                continue
            if r.y0 <= row_top + 25 and r.y1 >= row_top:
                area = max(0, min(r.x1, table_rect.x1) - max(r.x0, table_rect.x0))
                candidates.append((area, f))
        if candidates:
            _, best = max(candidates, key=lambda x: x[0])
            hdr_fill = f"{int(best[0]*255):02X}{int(best[1]*255):02X}{int(best[2]*255):02X}"
    except Exception:
        pass

    # Detect header text color from first row
    hdr_txt_color = None
    if first_row_spans:
        for spans in first_row_spans:
            for sp in spans:
                if sp.get("color", 0) != 0:
                    hdr_txt_color = sp["color"]
                    break
            if hdr_txt_color is not None:
                break

    return rows, hdr_fill, hdr_txt_color


def _build_word_table(document, rows, hdr_fill, hdr_txt_color, usable_cm=18.0):
    """Build an editable Word table from extracted row data."""
    from docx.shared import Cm, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT

    if not rows:
        return None

    ncols = max(len(r) for r in rows)
    rows = [r + [""] * (ncols - len(r)) for r in rows]

    # Column widths: give second col extra if 4 cols (SR#, Observation, Y/N, Remarks pattern)
    if ncols == 4:
        col_widths = [1.5, usable_cm - 1.5 - 2.0 - 5.0, 2.0, 5.0]
        col_widths[1] = max(4.0, col_widths[1])
    elif ncols == 3:
        col_widths = [1.5, usable_cm - 1.5 - 4.0, 4.0]
    else:
        each = usable_cm / ncols
        col_widths = [each] * ncols
        if ncols >= 3:
            bonus = min(3.0, usable_cm * 0.2)
            col_widths[1] += bonus
            share = bonus / (ncols - 1)
            for i in [0] + list(range(2, ncols)):
                col_widths[i] = max(1.0, col_widths[i] - share / (ncols-1))

    t = document.add_table(rows=len(rows), cols=ncols)
    t.alignment = WD_TABLE_ALIGNMENT.LEFT
    _set_table_borders(t, color="B0C4D8", sz="4")

    for ri, row_data in enumerate(rows):
        is_hdr = (ri == 0)
        for ci, txt in enumerate(row_data):
            cell = t.cell(ri, ci)
            try:
                cell.width = Cm(col_widths[ci])
            except Exception:
                pass
            _set_cell_shading(cell, hdr_fill if is_hdr else "FFFFFF")
            _set_cell_border(cell, "B0C4D8", "4")
            cell.text = ""
            p = cell.paragraphs[0]
            p.paragraph_format.space_before = Pt(2)
            p.paragraph_format.space_after = Pt(2)

            lines = (txt or "").split("\n")
            for li, line in enumerate(lines):
                if li > 0:
                    p.add_run().add_break()
                run = p.add_run(line)
                run.font.size = Pt(8 if is_hdr else 9)
                if is_hdr:
                    run.bold = True
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    if hdr_txt_color:
                        run.font.color.rgb = _pdf_rgb(hdr_txt_color)
                    else:
                        run.font.color.rgb = RGBColor(255, 255, 255)
                else:
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    # Bold SR# cells (e.g. "1.0", "2.0")
                    if ci == 0 and txt.strip() and "." in txt:
                        run.bold = True
    return t


@app.route("/api/pdf-to-word", methods=["POST"])
def pdf_to_word():
    """
    PDF → Word conversion.
    Strategy:
    - Detect table regions on each page.
    - Render non-table areas (header, footer, decorative graphics) as images.
    - Insert editable Word tables for the data portions.
    - This gives near pixel-perfect visual fidelity + editable cells.
    """
    from docx.shared import Cm, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT

    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a PDF file to convert."}), 400

    src = files[0]
    data = src.read()
    base_name = os.path.splitext(secure_filename(src.filename))[0] or "document"

    document = Document()
    for section in document.sections:
        section.page_width    = Cm(21)
        section.page_height   = Cm(29.7)
        section.left_margin   = Cm(0.5)
        section.right_margin  = Cm(0.5)
        section.top_margin    = Cm(0.3)
        section.bottom_margin = Cm(0.3)

    normal_style = document.styles["Normal"]
    normal_style.paragraph_format.space_after  = Pt(0)
    normal_style.paragraph_format.space_before = Pt(0)

    PAGE_W_CM = 20.0   # usable width after margins
    ZOOM = 3.0         # render quality

    with fitz.open(stream=data, filetype="pdf") as doc:
        for page_index, page in enumerate(doc):
            if page_index > 0:
                document.add_page_break()

            page_rect = page.rect   # full page bounds
            table_rects = _detect_table_regions(page)

            if not table_rects:
                # No tables found — render whole page as image (fallback)
                full_img = _clip_page_region(page, page_rect, zoom=2.5)
                _add_image_para(document, full_img, PAGE_W_CM)
                continue

            # Sort tables top-to-bottom
            table_rects.sort(key=lambda r: r.y0)

            # Build a list of regions: alternating image strips and table zones
            # Each item: ("image", y0, y1) or ("table", rect)
            regions = []
            cursor_y = page_rect.y0

            for tr in table_rects:
                # Gap before this table → image strip
                gap_y0 = cursor_y
                gap_y1 = tr.y0
                if gap_y1 > gap_y0 + 2:
                    regions.append(("image", gap_y0, gap_y1))
                regions.append(("table", tr))
                cursor_y = tr.y1

            # Remaining area after last table → image strip (footer)
            if cursor_y < page_rect.y1 - 2:
                regions.append(("image", cursor_y, page_rect.y1))

            # Render each region
            for region in regions:
                if region[0] == "image":
                    _, y0, y1 = region
                    clip = fitz.Rect(page_rect.x0, y0, page_rect.x1, y1)
                    img_bytes = _clip_page_region(page, clip, zoom=ZOOM)
                    _add_image_para(document, img_bytes, PAGE_W_CM)

                elif region[0] == "table":
                    tr = region[1]
                    rows, hdr_fill, hdr_txt_color = _extract_table_data(page, tr)
                    if rows:
                        sp = document.add_paragraph()
                        sp.paragraph_format.space_before = Pt(2)
                        sp.paragraph_format.space_after = Pt(2)
                        _build_word_table(document, rows, hdr_fill, hdr_txt_color, usable_cm=PAGE_W_CM)
                        sp2 = document.add_paragraph()
                        sp2.paragraph_format.space_before = Pt(2)
                        sp2.paragraph_format.space_after = Pt(2)
                    else:
                        # Table extraction failed — render as image
                        img_bytes = _clip_page_region(page, tr, zoom=ZOOM)
                        _add_image_para(document, img_bytes, PAGE_W_CM)

    buf = io.BytesIO()
    document.save(buf)
    buf.seek(0)

    return send_bytes(
        buf.getvalue(),
        f"{base_name}.docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )



# ---------------------------------------------------------------------------
# PDF -> PowerPoint (.pptx) - each page becomes a full-slide image
# ---------------------------------------------------------------------------

@app.route("/api/pdf-to-pptx", methods=["POST"])
def pdf_to_pptx():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a PDF file to convert."}), 400

    src = files[0]
    data = src.read()
    base_name = os.path.splitext(secure_filename(src.filename))[0] or "presentation"

    zoom = 2.0
    mat = fitz.Matrix(zoom, zoom)

    prs = Presentation()

    with fitz.open(stream=data, filetype="pdf") as doc:
        if len(doc) == 0:
            return jsonify({"error": "The PDF has no pages."}), 400

        # Size the slides to match the first page's aspect ratio
        first_rect = doc[0].rect
        prs.slide_width = Emu(int(first_rect.width * 12700))
        prs.slide_height = Emu(int(first_rect.height * 12700))
        blank_layout = prs.slide_layouts[6]

        for page in doc:
            pix = page.get_pixmap(matrix=mat)
            img_bytes = pix.tobytes("png")

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                Emu(0),
                Emu(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    return send_bytes(
        buf.getvalue(),
        f"{base_name}.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


# ---------------------------------------------------------------------------
# PDF -> Excel (.xlsx) - extracts detected tables, falls back to raw text lines
# ---------------------------------------------------------------------------

@app.route("/api/pdf-to-excel", methods=["POST"])
def pdf_to_excel():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a PDF file to convert."}), 400

    src = files[0]
    data = src.read()
    base_name = os.path.splitext(secure_filename(src.filename))[0] or "spreadsheet"

    wb = Workbook()
    wb.remove(wb.active)  # start with no sheets, add one per page below

    with fitz.open(stream=data, filetype="pdf") as doc:
        for page_index, page in enumerate(doc):
            sheet_name = f"Page {page_index + 1}"[:31]
            ws = wb.create_sheet(title=sheet_name)

            wrote_table = False
            try:
                tables = page.find_tables()
                if tables and tables.tables:
                    row_cursor = 1
                    for table in tables.tables:
                        extracted = table.extract()
                        for row in extracted:
                            for col_index, cell_value in enumerate(row, start=1):
                                ws.cell(
                                    row=row_cursor,
                                    column=col_index,
                                    value=cell_value if cell_value is not None else "",
                                )
                            row_cursor += 1
                        row_cursor += 1  # blank row between tables
                        wrote_table = True
            except Exception:
                wrote_table = False

            if not wrote_table:
                # Fall back to dumping each line of text into column A
                text = page.get_text()
                for row_index, line in enumerate(text.split("\n"), start=1):
                    if line.strip():
                        ws.cell(row=row_index, column=1, value=line)

    if not wb.sheetnames:
        wb.create_sheet(title="Sheet1")

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)

    return send_bytes(
        buf.getvalue(),
        f"{base_name}.xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )


# ---------------------------------------------------------------------------
# Remove pages
# ---------------------------------------------------------------------------

@app.route("/api/remove-pages", methods=["POST"])
def remove_pages():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a PDF file."}), 400
    pages_spec = request.form.get("pages", "")
    if not pages_spec.strip():
        return jsonify({"error": "Tell me which pages to remove, e.g. 2,4-5"}), 400

    src = files[0]
    data = src.read()

    with fitz.open(stream=data, filetype="pdf") as doc:
        to_remove = sorted(set(parse_page_spec(pages_spec, len(doc))), reverse=True)
        if not to_remove:
            return jsonify({"error": "No valid page numbers found in that range."}), 400
        for idx in to_remove:
            doc.delete_page(idx)
        if len(doc) == 0:
            return jsonify({"error": "That would remove every page. Leave at least one."}), 400
        out_bytes = doc.tobytes()

    return send_bytes(out_bytes, "pages_removed.pdf", "application/pdf")


# ---------------------------------------------------------------------------
# Extract / reorder pages ("Organize PDF")
# ---------------------------------------------------------------------------

@app.route("/api/extract-pages", methods=["POST"])
def extract_pages():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a PDF file."}), 400
    pages_spec = request.form.get("pages", "")
    if not pages_spec.strip():
        return jsonify({"error": "Tell me which pages to keep and in what order, e.g. 3,1,2"}), 400

    src = files[0]
    data = src.read()

    with fitz.open(stream=data, filetype="pdf") as doc:
        page_order = parse_page_spec(pages_spec, len(doc))
        if not page_order:
            return jsonify({"error": "No valid page numbers found in that range."}), 400

        new_doc = fitz.open()
        for idx in page_order:
            new_doc.insert_pdf(doc, from_page=idx, to_page=idx)
        out_bytes = new_doc.tobytes()
        new_doc.close()

    return send_bytes(out_bytes, "organized.pdf", "application/pdf")


# ---------------------------------------------------------------------------
# Add page numbers
# ---------------------------------------------------------------------------

@app.route("/api/add-page-numbers", methods=["POST"])
def add_page_numbers():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a PDF file."}), 400

    try:
        start = int(request.form.get("start", 1))
    except ValueError:
        start = 1
    position = request.form.get("position", "bottom-center")

    src = files[0]
    data = src.read()

    with fitz.open(stream=data, filetype="pdf") as doc:
        for i, page in enumerate(doc):
            label = str(start + i)
            rect = page.rect
            margin = 24
            if position == "bottom-left":
                point = fitz.Point(margin, rect.height - margin)
            elif position == "bottom-right":
                point = fitz.Point(rect.width - margin - 20, rect.height - margin)
            elif position == "top-center":
                point = fitz.Point(rect.width / 2 - 8, margin)
            else:  # bottom-center
                point = fitz.Point(rect.width / 2 - 8, rect.height - margin)

            page.insert_text(point, label, fontsize=11, color=(0, 0, 0))

        out_bytes = doc.tobytes()

    return send_bytes(out_bytes, "numbered.pdf", "application/pdf")


# ---------------------------------------------------------------------------
# Watermark PDF (diagonal repeated text)
# ---------------------------------------------------------------------------

@app.route("/api/watermark", methods=["POST"])
def watermark_pdf():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a PDF file."}), 400

    text = request.form.get("text", "CONFIDENTIAL").strip() or "CONFIDENTIAL"
    try:
        opacity = float(request.form.get("opacity", 0.3))
    except ValueError:
        opacity = 0.3
    opacity = max(0.05, min(opacity, 1.0))

    src = files[0]
    data = src.read()

    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            rect = page.rect
            page.insert_text(
                fitz.Point(rect.width * 0.15, rect.height * 0.55),
                text,
                fontsize=max(24, int(rect.width / 12)),
                rotate=45,
                color=(0.6, 0.6, 0.6),
                fill_opacity=opacity,
                overlay=True,
            )
        out_bytes = doc.tobytes()

    return send_bytes(out_bytes, "watermarked.pdf", "application/pdf")


# ---------------------------------------------------------------------------
# Protect PDF (add a password)
# ---------------------------------------------------------------------------

@app.route("/api/protect", methods=["POST"])
def protect_pdf():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a PDF file."}), 400
    password = request.form.get("password", "").strip()
    if not password:
        return jsonify({"error": "Enter a password to protect the PDF with."}), 400

    src = files[0]
    data = src.read()

    with fitz.open(stream=data, filetype="pdf") as doc:
        out_bytes = doc.tobytes(
            encryption=fitz.PDF_ENCRYPT_AES_256,
            owner_pw=password,
            user_pw=password,
            permissions=int(
                fitz.PDF_PERM_PRINT
                | fitz.PDF_PERM_COPY
                | fitz.PDF_PERM_ANNOTATE
            ),
        )

    return send_bytes(out_bytes, "protected.pdf", "application/pdf")


# ---------------------------------------------------------------------------
# Unlock PDF (remove a known password)
# ---------------------------------------------------------------------------

@app.route("/api/unlock", methods=["POST"])
def unlock_pdf():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a PDF file."}), 400
    password = request.form.get("password", "").strip()

    src = files[0]
    data = src.read()

    doc = fitz.open(stream=data, filetype="pdf")
    try:
        if doc.needs_pass:
            if not password:
                return jsonify({"error": "This PDF is password-protected. Enter the password."}), 400
            if not doc.authenticate(password):
                return jsonify({"error": "That password didn't work."}), 400
        out_bytes = doc.tobytes()
    finally:
        doc.close()

    return send_bytes(out_bytes, "unlocked.pdf", "application/pdf")


# ---------------------------------------------------------------------------
# Crop PDF (trim a margin off every page)
# ---------------------------------------------------------------------------

@app.route("/api/crop", methods=["POST"])
def crop_pdf():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a PDF file."}), 400

    try:
        margin = float(request.form.get("margin", 36))
    except ValueError:
        margin = 36
    margin = max(0, margin)

    src = files[0]
    data = src.read()

    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            rect = page.rect
            new_rect = fitz.Rect(
                rect.x0 + margin,
                rect.y0 + margin,
                rect.x1 - margin,
                rect.y1 - margin,
            )
            if new_rect.width > 10 and new_rect.height > 10:
                page.set_cropbox(new_rect)
        out_bytes = doc.tobytes()

    return send_bytes(out_bytes, "cropped.pdf", "application/pdf")


# ---------------------------------------------------------------------------
# Word (.docx) -> PDF  (text-based conversion, not full layout fidelity)
# ---------------------------------------------------------------------------

@app.route("/api/word-to-pdf", methods=["POST"])
def word_to_pdf():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a .docx file."}), 400

    src = files[0]
    base_name = os.path.splitext(secure_filename(src.filename))[0] or "document"
    document = Document(src.stream)

    buf = io.BytesIO()
    pdf_doc = SimpleDocTemplate(buf, pagesize=letter)
    styles = getSampleStyleSheet()
    flowables = []

    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            flowables.append(Spacer(1, 10))
            continue
        style_name = "Heading2" if para.style and "Heading" in (para.style.name or "") else "Normal"
        flowables.append(Paragraph(text.replace("&", "&amp;").replace("<", "&lt;"), styles[style_name]))
        flowables.append(Spacer(1, 6))

    if not flowables:
        flowables = [Paragraph("(No text content found in this document.)", styles["Normal"])]

    pdf_doc.build(flowables)
    buf.seek(0)

    return send_bytes(buf.getvalue(), f"{base_name}.pdf", "application/pdf")


# ---------------------------------------------------------------------------
# PowerPoint (.pptx) -> PDF (one page per slide, text content only)
# ---------------------------------------------------------------------------

@app.route("/api/pptx-to-pdf", methods=["POST"])
def pptx_to_pdf():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a .pptx file."}), 400

    src = files[0]
    base_name = os.path.splitext(secure_filename(src.filename))[0] or "presentation"
    prs = Presentation(src.stream)

    buf = io.BytesIO()
    c = pdf_canvas.Canvas(buf, pagesize=letter)
    width, height = letter

    for slide in prs.slides:
        y = height - 60
        c.setFont("Helvetica-Bold", 16)
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text.strip())

        if texts:
            c.drawString(50, y, texts[0][:90])
            y -= 40
            c.setFont("Helvetica", 12)
            for block in texts[1:]:
                for line in block.split("\n"):
                    if y < 50:
                        c.showPage()
                        y = height - 60
                        c.setFont("Helvetica", 12)
                    c.drawString(60, y, line[:100])
                    y -= 18
        else:
            c.setFont("Helvetica", 12)
            c.drawString(50, y, "(No text content on this slide.)")

        c.showPage()

    c.save()
    buf.seek(0)

    return send_bytes(buf.getvalue(), f"{base_name}.pdf", "application/pdf")


# ---------------------------------------------------------------------------
# Excel (.xlsx) -> PDF (one table per sheet)
# ---------------------------------------------------------------------------

@app.route("/api/excel-to-pdf", methods=["POST"])
def excel_to_pdf():
    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload a .xlsx file."}), 400

    src = files[0]
    base_name = os.path.splitext(secure_filename(src.filename))[0] or "spreadsheet"
    wb = load_workbook(src.stream, data_only=True)

    buf = io.BytesIO()
    pdf_doc = SimpleDocTemplate(buf, pagesize=letter)
    styles = getSampleStyleSheet()
    flowables = []

    for sheet in wb.worksheets:
        flowables.append(Paragraph(sheet.title, styles["Heading2"]))
        flowables.append(Spacer(1, 8))

        rows = []
        for row in sheet.iter_rows(values_only=True):
            rows.append(["" if v is None else str(v) for v in row])

        if rows:
            # Cap columns so wide sheets don't overflow the page unreadably
            max_cols = 10
            rows = [r[:max_cols] for r in rows]
            table = Table(rows, repeatRows=1)
            table.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("BACKGROUND", (0, 0), (-1, 0), colors.whitesmoke),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
            ]))
            flowables.append(table)
        else:
            flowables.append(Paragraph("(Empty sheet)", styles["Normal"]))

        flowables.append(Spacer(1, 20))

    pdf_doc.build(flowables)
    buf.seek(0)

    return send_bytes(buf.getvalue(), f"{base_name}.pdf", "application/pdf")


# ---------------------------------------------------------------------------
# Image → OCR Text / Word (.docx)
# Accepts JPG / PNG / BMP / TIFF / WebP.  Requires Tesseract on the server
# (installed via nixpacks.toml for Railway deployments).
# ---------------------------------------------------------------------------

def _ocr_image(pil_img, lang="eng"):
    """
    Run Tesseract OCR on a PIL image.
    Raises RuntimeError with a user-friendly message if Tesseract is missing.
    """
    import shutil
    if not shutil.which("tesseract"):
        raise RuntimeError(
            "Tesseract OCR is not installed on this server. "
            "Add nixpacks.toml with tesseract to your Railway project and redeploy."
        )
    try:
        import pytesseract
    except ImportError:
        raise RuntimeError(
            "pytesseract is missing. Add it to requirements.txt and redeploy."
        )
    # Upscale tiny images — Tesseract is much more accurate at ≥ ~200 DPI
    w, h = pil_img.size
    if max(w, h) < 1200:
        scale = 1200 / max(w, h)
        pil_img = pil_img.resize(
            (int(w * scale), int(h * scale)), Image.LANCZOS
        )
    return pytesseract.image_to_string(pil_img, lang=lang, config="--psm 3")


@app.route("/api/image-to-ocr", methods=["POST"])
def image_to_ocr():
    import re
    from docx.shared import Cm, RGBColor

    files = get_uploaded_files()
    if not files:
        return jsonify({"error": "Upload one or more image files."}), 400

    output_format = request.form.get("output_format", "txt")   # "txt" or "docx"
    lang          = request.form.get("lang", "eng")            # tesseract lang code

    # Whitelist: only letters + '+' (e.g. "eng", "ara", "ara+eng")
    if not re.match(r'^[a-zA-Z+]{2,20}$', lang):
        lang = "eng"

    # ── OCR every uploaded image ─────────────────────────────────────────────
    page_texts = []
    for f in files:
        name = os.path.splitext(secure_filename(f.filename))[0] or "image"
        try:
            img = Image.open(f.stream)
            # Tesseract wants RGB or greyscale — flatten RGBA onto white bg
            if img.mode == "RGBA":
                bg = Image.new("RGB", img.size, (255, 255, 255))
                bg.paste(img, mask=img.split()[3])
                img = bg
            elif img.mode not in ("RGB", "L"):
                img = img.convert("RGB")

            ocr_text = _ocr_image(img, lang=lang)
            page_texts.append((name, ocr_text))

        except RuntimeError as e:
            # Tesseract not installed — return a clear 503 immediately
            return jsonify({"error": str(e)}), 503
        except Exception as e:
            page_texts.append((name, f"[Could not read this image: {e}]"))

    if not page_texts:
        return jsonify({"error": "No images could be processed."}), 400

    base_name = page_texts[0][0] or "ocr_result"

    # ── Plain-text output ────────────────────────────────────────────────────
    if output_format == "txt":
        parts = []
        for name, text in page_texts:
            if len(page_texts) > 1:
                parts.append(f"=== {name} ===")
            parts.append(text.strip())
            parts.append("")
        raw = "\n".join(parts).strip().encode("utf-8")
        return send_bytes(raw, f"{base_name}_ocr.txt", "text/plain; charset=utf-8")

    # ── Word (.docx) output ──────────────────────────────────────────────────
    document = Document()
    for section in document.sections:
        section.page_width    = Cm(21)
        section.page_height   = Cm(29.7)
        section.left_margin   = Cm(2.0)
        section.right_margin  = Cm(2.0)
        section.top_margin    = Cm(2.0)
        section.bottom_margin = Cm(2.0)

    normal = document.styles["Normal"]
    normal.paragraph_format.space_after  = Pt(2)
    normal.paragraph_format.space_before = Pt(0)

    for idx, (name, text) in enumerate(page_texts):
        if idx > 0:
            document.add_page_break()
        if len(page_texts) > 1:
            document.add_paragraph(style="Heading 1").add_run(name)
        for line in text.splitlines():
            p = document.add_paragraph()
            p.paragraph_format.space_after  = Pt(0)
            p.paragraph_format.space_before = Pt(0)
            run = p.add_run(line)
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(0, 0, 0)

    buf = io.BytesIO()
    document.save(buf)
    buf.seek(0)
    return send_bytes(
        buf.getvalue(),
        f"{base_name}_ocr.docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )


# ---------------------------------------------------------------------------
# Error handlers
# ---------------------------------------------------------------------------

@app.errorhandler(413)
def too_large(_e):
    return jsonify({"error": "File too large. Max size is 50MB."}), 413


@app.errorhandler(500)
def server_error(e):
    return jsonify({"error": f"Something went wrong: {str(e)}"}), 500


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 7860))
    app.run(host="0.0.0.0", port=port, debug=False)